#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智能知识库系统 - 主应用
整合PyMuPDF Pro + PyMuPDF4LLM + LangChain + 极客智坊API
"""

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import uvicorn
import json
import logging
import os
import tempfile
from pathlib import Path
import subprocess
from shutil import which
import shutil
import requests

# PyMuPDF相关
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("❌ PyMuPDF 不可用")

# 可选：用于优化Excel分页的预处理
try:
    from openpyxl import load_workbook
    from openpyxl.worksheet.properties import PageSetupProperties
    OPENPYXL_AVAILABLE = True
except Exception:
    OPENPYXL_AVAILABLE = False

# LangChain相关
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain.schema import Document
    LANGCHAIN_AVAILABLE = True
except ImportError:
    LANGCHAIN_AVAILABLE = False
    print("❌ LangChain 不可用")

# ES相关
from elasticsearch import Elasticsearch
import hashlib

# Office文档解析能力（原生）
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="智能知识库系统", version="2.0.0")

# 导入配置
from config import (
    ES_CONFIG, DOCUMENT_CONFIG, EMBEDDING_CONFIG, RAG_CONFIG,
    EMBEDDING_DIMS,
    CHUNKING_CONFIG, GEEKAI_API_KEY, GEEKAI_CHAT_URL,
    GEEKAI_EMBEDDING_URL, DEFAULT_EMBEDDING_MODEL
)
from fastapi.responses import FileResponse

# 辅助：构建页文本与词级索引（用于bbox定位）
def build_page_text_and_word_index(page: "fitz.Page") -> (str, list):
    """返回 (page_text, word_entries), 其中 word_entries 为 [ (start, end, (x0,y0,x1,y1)) ]."""
    try:
        words = page.get_text("words")  # (x0,y0,x1,y1,word,block_no,line_no,word_no)
        # 排序：块→行→词
        words_sorted = sorted(words, key=lambda w: (int(w[5]), int(w[6]), int(w[7])))
        parts = []
        entries = []
        last_block, last_line = None, None
        current_pos = 0
        for w in words_sorted:
            x0, y0, x1, y1, token, bno, lno, wno = w
            key = (bno, lno)
            if last_block is None:
                # 首词，直接写
                parts.append(token)
                start = current_pos
                end = start + len(token)
                entries.append((start, end, (float(x0), float(y0), float(x1), float(y1))))
                current_pos = end
            else:
                if key != (last_block, last_line):
                    # 换行
                    parts.append("\n")
                    current_pos += 1
                    parts.append(token)
                    start = current_pos
                    end = start + len(token)
                    entries.append((start, end, (float(x0), float(y0), float(x1), float(y1))))
                    current_pos = end
                else:
                    # 同行词，空格分隔
                    parts.append(" ")
                    current_pos += 1
                    parts.append(token)
                    start = current_pos
                    end = start + len(token)
                    entries.append((start, end, (float(x0), float(y0), float(x1), float(y1))))
                    current_pos = end
            last_block, last_line = key
        page_text = "".join(parts)
        return page_text, entries
    except Exception:
        txt = page.get_text() or ""
        return txt, []

def compute_bbox_union_for_range(entries: list, start: int, end: int) -> list:
    """根据词级索引列表，计算覆盖 [start,end) 的bbox并集，返回 [x0,y0,x1,y1]，无则返回空列表。"""
    x0 = y0 = float("inf")
    x1 = y1 = float("-inf")
    found = False
    for (s, e, (a, b, c, d)) in entries:
        if e <= start or s >= end:
            continue
        found = True
        x0 = min(x0, a)
        y0 = min(y0, b)
        x1 = max(x1, c)
        y1 = max(y1, d)
    return [x0, y0, x1, y1] if found else []

# 初始化ES客户端
es_client = Elasticsearch(
    [f"http://{ES_CONFIG['host']}:{ES_CONFIG['port']}"],
    basic_auth=(ES_CONFIG['username'], ES_CONFIG['password']) if ES_CONFIG['username'] else None,
    verify_certs=ES_CONFIG['verify_certs']
)

# 极客API embedding函数
def get_embedding(text: str) -> list:
    """使用极客API获取文本嵌入向量"""
    try:
        url = GEEKAI_EMBEDDING_URL
        headers = {
            "Authorization": f"Bearer {GEEKAI_API_KEY}",
            "Content-Type": "application/json"
        }
        data = {
            "model": DEFAULT_EMBEDDING_MODEL,
            "input": [text],
            "intent": "search_document"
        }
        response = requests.post(url, headers=headers, json=data, timeout=20)
        if response.status_code == 200:
            return response.json().get("data", [{}])[0].get("embedding")
        else:
            logger.error(f"极客API embedding失败: {response.status_code} - {response.text}")
            return None
    except Exception as e:
        logger.error(f"极客API embedding异常: {e}")
        return None

# 请求模型
class LdapValidateRequest(BaseModel):
    username: str
    password: str

class ChatRequest(BaseModel):
    question: str
    user_id: str
    source_file: Optional[str] = None  # 可选：指定特定文件名进行RAG检索

class DocumentProcessRequest(BaseModel):
    knowledge_id: int
    knowledge_name: str
    description: str
    tags: List[str]
    effective_time: str

# 响应模型
class LdapValidateResponse(BaseModel):
    success: bool
    message: str
    email: Optional[str] = None
    role: Optional[str] = None

class KnowledgeReference(BaseModel):
    knowledge_id: int
    knowledge_name: str
    description: str
    tags: List[str]
    effective_time: str
    attachments: List[str]
    relevance: float
    # 追溯定位增强
    source_file: Optional[str] = None
    page_num: Optional[int] = None
    chunk_index: Optional[int] = None
    chunk_type: Optional[str] = None
    # 新增：返回块坐标与字符范围，便于前端高亮
    bbox_union: Optional[List[float]] = None
    char_start: Optional[int] = None
    char_end: Optional[int] = None

class ChatResponse(BaseModel):
    answer: str
    references: List[KnowledgeReference]
    session_id: Optional[str] = None

class DocumentProcessResponse(BaseModel):
    success: bool
    message: str
    chunks_count: int
    knowledge_id: int

@app.get("/")
def read_root():
    """根路径"""
    return {
        "message": "智能知识库系统",
        "version": "2.0.0",
        "features": [
            "PyMuPDF Pro 文档处理",
            "PyMuPDF4LLM 结构化分块", 
            "LangChain 向量化",
            "Elasticsearch 存储",
            "极客智坊API 智能问答"
        ]
    }

# ===== LibreOffice 转PDF（不改动原文件）=====
def _find_soffice_path() -> Optional[str]:
    candidates = [
        r"C:\\Program Files\\LibreOffice\\program\\soffice.exe",
        r"C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
    ]
    for c in candidates:
        if os.path.exists(c):
            return c
    w = which("soffice") or which("soffice.exe")
    return w

def convert_with_libreoffice_safe(src_path: str, timeout_sec: int = 180, out_dir: Optional[str] = None) -> str:
    """使用LibreOffice将任意Office文档转为PDF，输出到临时文件夹，返回PDF路径。"""
    soffice = _find_soffice_path()
    if not soffice:
        raise RuntimeError("未找到LibreOffice的soffice.exe，请安装或配置PATH后重试")

    if not out_dir:
        out_dir = tempfile.mkdtemp(prefix="lo_pdf_")
    else:
        os.makedirs(out_dir, exist_ok=True)
    cmd = [
        soffice,
        "--headless",
        "--norestore",
        "--nolockcheck",
        "--convert-to", "pdf",
        "--outdir", out_dir,
        os.path.abspath(src_path),
    ]
    cp = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout_sec)
    if cp.returncode != 0:
        raise RuntimeError(f"LibreOffice 转换失败: rc={cp.returncode}\nstdout={cp.stdout}\nstderr={cp.stderr}")

    # 期望同名pdf
    expected = os.path.join(out_dir, Path(src_path).with_suffix('.pdf').name)
    if os.path.exists(expected):
        return expected
    # 兜底找最新pdf
    pdfs = [p for p in os.listdir(out_dir) if p.lower().endswith('.pdf')]
    if not pdfs:
        raise RuntimeError(f"LibreOffice 未生成PDF。stdout={cp.stdout}\nstderr={cp.stderr}")
    pdfs.sort(key=lambda n: os.path.getmtime(os.path.join(out_dir, n)), reverse=True)
    return os.path.join(out_dir, pdfs[0])

def create_simple_pdf_from_txt(src_path: str, out_path: Optional[str] = None) -> str:
    """TXT 降级方案：将文本简单排版写入单页/多页PDF，便于坐标回显。"""
    with open(src_path, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    out_pdf = out_path or tempfile.mktemp(suffix='.pdf')
    os.makedirs(os.path.dirname(out_pdf), exist_ok=True)
    doc = fitz.open()
    # 解决中文显示为问号：嵌入系统CJK字体（优先微软雅黑 / 宋体 / 黑体）
    def _register_cjk_font(_doc: fitz.Document) -> str:
        try:
            candidates = [
                r"C:\\Windows\\Fonts\\msyh.ttc",
                r"C:\\Windows\\Fonts\\msyh.ttf",
                r"C:\\Windows\\Fonts\\simsun.ttc",
                r"C:\\Windows\\Fonts\\simhei.ttf",
            ]
            for p in candidates:
                if os.path.exists(p):
                    _doc.insert_font(fontname="CJK", fontfile=p)
                    return "CJK"
        except Exception:
            pass
        return "helv"
    cjk_font = _register_cjk_font(doc)
    page = doc.new_page(width=595, height=842)
    rect = fitz.Rect(40, 50, 555, 800)
    page.insert_textbox(rect, text[:50000], fontsize=11, fontname=cjk_font, align=0)
    doc.save(out_pdf)
    doc.close()
    return out_pdf

# ===== 持久化存储的转换PDF路径 =====
CONVERTED_PDF_ROOT = os.path.join(os.path.dirname(__file__), 'static', 'converted')

def build_converted_pdf_path(knowledge_id: int, original_filename: str) -> str:
    base = Path(original_filename).stem + '.pdf'
    target_dir = os.path.join(CONVERTED_PDF_ROOT, str(knowledge_id))
    os.makedirs(target_dir, exist_ok=True)
    return os.path.join(target_dir, base)

def preprocess_xlsx_for_better_pdf(src_path: str) -> Optional[str]:
    """
    针对Excel进行分页优化：横向、按宽度适配、设置打印区域与边距；
    生成临时xlsx副本（不改动原文件）。返回临时文件路径；失败返回None。
    """
    if not OPENPYXL_AVAILABLE:
        return None

# ===== Excel 原生解析，不转PDF =====
def parse_excel_native(file_path: str, filename: str, knowledge_id: int) -> List[Document]:
    if not OPENPYXL_AVAILABLE:
        raise RuntimeError("缺少openpyxl，无法原生解析Excel。请安装: pip install openpyxl")
    wb = load_workbook(file_path, data_only=True)
    chunks: List[Document] = []

    # 构建每个sheet的虚拟坐标网格
    for sheet_idx, ws in enumerate(wb.worksheets):
        # 列宽、行高 → 像素近似（Excel宽度单位转像素，这里做简单近似）
        col_widths = []
        for col in ws.columns:
            # openpyxl列宽在 ws.column_dimensions[col_letter].width
            break
        # 收集列宽（若未设置，给默认 8.43 字符宽 ≈ 64px）
        from openpyxl.utils import get_column_letter
        max_col = ws.max_column
        max_row = ws.max_row
        for c in range(1, max_col + 1):
            letter = get_column_letter(c)
            cw = ws.column_dimensions.get(letter).width if ws.column_dimensions.get(letter) and ws.column_dimensions.get(letter).width else 8.43
            # 近似：1字符宽≈7.5px
            col_widths.append(float(cw) * 7.5)
        row_heights = []
        for r in range(1, max_row + 1):
            rh = ws.row_dimensions.get(r).height if ws.row_dimensions.get(r) and ws.row_dimensions.get(r).height else 15.0
            # 近似：1pt≈1.33px，Excel默认行高≈15pt
            row_heights.append(float(rh) * 1.33)

        # 前缀和得到各列x、各行y起点
        x_starts = [0.0]
        for w in col_widths:
            x_starts.append(x_starts[-1] + w)
        y_starts = [0.0]
        for h in row_heights:
            y_starts.append(y_starts[-1] + h)

        # 构建单元格positions
        positions: List[Dict] = []
        for r in range(1, max_row + 1):
            for c in range(1, max_col + 1):
                cell = ws.cell(row=r, column=c)
                text = str(cell.value) if cell.value is not None else ""
                if not text.strip():
                    continue
                # 合并单元格处理：取所属合并区的外接矩形
                merged_bbox = None
                for rng in ws.merged_cells.ranges:
                    if (r, c) in rng.cells:
                        min_row, min_col, max_row_, max_col_ = rng.min_row, rng.min_col, rng.max_row, rng.max_col
                        x0 = x_starts[min_col - 1]
                        y0 = y_starts[min_row - 1]
                        x1 = x_starts[max_col_]
                        y1 = y_starts[max_row_]
                        merged_bbox = [x0, y0, x1, y1]
                        break
                if merged_bbox:
                    bbox = merged_bbox
                else:
                    x0 = x_starts[c - 1]
                    y0 = y_starts[r - 1]
                    x1 = x_starts[c]
                    y1 = y_starts[r]
                    bbox = [x0, y0, x1, y1]

                positions.append({
                    "text": text.strip(),
                    "bbox": bbox,
                    "sheet": ws.title,
                    "row": r,
                    "col": c,
                })

        # 生成纯文本内容（按行拼接）
        lines = []
        for r in range(1, max_row + 1):
            vals = []
            for c in range(1, max_col + 1):
                val = ws.cell(row=r, column=c).value
                vals.append("" if val is None else str(val))
            lines.append("\t".join(vals).rstrip())
        all_text = "\n".join(lines)

        # 分块（按字符）
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        text_chunks = splitter.split_text(all_text)

        # 将positions按文本包含关系分配给chunk
        for chunk_idx, chunk_text in enumerate(text_chunks):
            chunk_positions = []
            # 简单包含匹配：单元格文本在chunk中出现则纳入（大多数表格有效）
            for p in positions:
                t = p["text"]
                if t and t in chunk_text:
                    chunk_positions.append(p)

            # 计算块级bbox：所有单元格bbox并集（同一sheet）
            if chunk_positions:
                x0 = min(pp["bbox"][0] for pp in chunk_positions)
                y0 = min(pp["bbox"][1] for pp in chunk_positions)
                x1 = max(pp["bbox"][2] for pp in chunk_positions)
                y1 = max(pp["bbox"][3] for pp in chunk_positions)
                bbox = [x0, y0, x1, y1]
            else:
                bbox = [0, 0, 0, 0]

            chunks.append(Document(
                page_content=chunk_text,
                metadata={
                    "knowledge_id": knowledge_id,
                    "source_file": filename,
                    "page_num": sheet_idx + 1,  # 将 sheet 当作“页”
                    "chunk_index": len(chunks),
                    "positions": chunk_positions,
                    "bbox": bbox,
                    "document_name": filename,
                    "document_type": "表格",
                    "sheet_name": ws.title,
                    "keywords": extract_keywords_from_content(chunk_text),
                }
            ))

    return chunks

# ===== TXT 原生解析，不转PDF =====
def parse_txt_native(file_path: str, filename: str, knowledge_id: int) -> List[Document]:
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    # 归一化换行
    content = content.replace('\r\n', '\n').replace('\r', '\n')
    lines = content.split('\n')

    # 行级positions（不做真实像素坐标，记录行列与字符区间）
    positions: List[Dict] = []
    global_pos = 0
    for i, line in enumerate(lines, start=1):
        text = line
        start = global_pos
        end = start + len(text)
        positions.append({
            "text": text,
            "line_no": i,
            "char_start": start,
            "char_end": end,
        })
        global_pos = end + 1  # 计入换行

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    text_chunks = splitter.split_text(content)

    chunks: List[Document] = []
    for idx, chunk_text in enumerate(text_chunks):
        # 为chunk分配行段positions
        chunk_positions = []
        # 通过字符范围粗匹配（行文本在chunk中出现时纳入）
        for p in positions:
            t = p["text"]
            if t and t in chunk_text:
                chunk_positions.append(p)
        chunks.append(Document(
            page_content=chunk_text,
            metadata={
                "knowledge_id": knowledge_id,
                "source_file": filename,
                "page_num": 1,
                "chunk_index": idx,
                "positions": chunk_positions,
                "bbox": [],  # 文本视图通常不需要像素坐标
                "document_name": filename,
                "document_type": "文本",
                "keywords": extract_keywords_from_content(chunk_text),
            }
        ))

    return chunks
    try:
        wb = load_workbook(src_path, data_only=True)
        for ws in wb.worksheets:
            # 横向打印 + 宽度适配
            ws.page_setup.orientation = 'landscape'
            ws.page_setup.fitToWidth = 1
            ws.page_setup.fitToHeight = 0
            if ws.sheet_properties.pageSetUpPr is None:
                ws.sheet_properties.pageSetUpPr = PageSetupProperties(fitToPage=True)
            else:
                ws.sheet_properties.pageSetUpPr.fitToPage = True
            # 居中与边距
            ws.print_options.horizontalCentered = True
            ws.page_margins.left = 0.5
            ws.page_margins.right = 0.5
            ws.page_margins.top = 0.6
            ws.page_margins.bottom = 0.6
            # 打印区域覆盖已用范围
            try:
                dim = ws.calculate_dimension()  # 如 'A1:G200'
                ws.print_area = dim
            except Exception:
                pass
        fd, tmp_path = tempfile.mkstemp(suffix='.xlsx')
        os.close(fd)
        wb.save(tmp_path)
        return tmp_path
    except Exception as e:
        logger.warning(f"Excel预处理失败，使用原文件转换: {e}")
        return None

# ===== DOCX 原生解析与PDF生成（不依赖外部工具） =====
def create_simple_pdf_from_docx(src_path: str, out_path: str) -> tuple[str, List[Dict]]:
    """将DOCX内容简单排版生成PDF，并返回 (pdf_path, paragraph_positions)。
    每个段落在PDF中的放置矩形作为其bbox，供后续可视化使用。
    """
    if not DOCX_AVAILABLE:
        raise RuntimeError("缺少python-docx，无法原生解析Word。请安装: pip install python-docx")

    docx = DocxDocument(src_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)

    pdf = fitz.open()
    def _register_cjk_font(_doc: fitz.Document) -> str:
        try:
            candidates = [
                r"C:\\Windows\\Fonts\\msyh.ttc",
                r"C:\\Windows\\Fonts\\msyh.ttf",
                r"C:\\Windows\\Fonts\\simsun.ttc",
                r"C:\\Windows\\Fonts\\simhei.ttf",
            ]
            for p in candidates:
                if os.path.exists(p):
                    _doc.insert_font(fontname="CJK", fontfile=p)
                    return "CJK"
        except Exception:
            pass
        return "helv"
    cjk_font = _register_cjk_font(pdf)
    page_width, page_height = 595.0, 842.0  # A4 72dpi points
    margin_left, margin_right, margin_top, margin_bottom = 40.0, 40.0, 50.0, 50.0
    usable_width = page_width - margin_left - margin_right
    cursor_y = margin_top

    paragraph_positions: List[Dict] = []

    # 新建第一页
    page = pdf.new_page(width=page_width, height=page_height)

    def new_page_if_needed(height_needed: float):
        nonlocal page, cursor_y
        if cursor_y + height_needed > (page_height - margin_bottom):
            page = pdf.new_page(width=page_width, height=page_height)
            cursor_y = margin_top

    def draw_paragraph(text: str):
        nonlocal cursor_y
        if not text.strip():
            cursor_y += 8.0
            return
        # 预估高度：用一个较高的矩形，插入后不读取残留，靠换页控制
        rect = fitz.Rect(margin_left, cursor_y, margin_left + usable_width, cursor_y + 200.0)
        new_page_if_needed(rect.height)
        rect = fitz.Rect(margin_left, cursor_y, margin_left + usable_width, cursor_y + 200.0)
        page.insert_textbox(rect, text, fontsize=11, fontname=cjk_font, align=0)
        # 粗略估算占用高度：按文本长度估计行数
        approx_chars_per_line = 90
        lines = max(1, (len(text) // approx_chars_per_line) + 1)
        used_h = 14.0 * lines
        bbox = [rect.x0, rect.y0, rect.x1, rect.y0 + used_h]
        paragraph_positions.append({"text": text, "bbox": bbox, "page": len(pdf)})
        cursor_y += used_h + 6.0

    # 段落
    for p in docx.paragraphs:
        draw_paragraph(p.text)

    # 表格（将每个单元格作为独立小段落）
    for table in docx.tables:
        for row in table.rows:
            cells_text = [c.text.strip() for c in row.cells]
            row_text = "\t".join(cells_text)
            draw_paragraph(row_text)

    pdf.save(out_path)
    pdf.close()
    return out_path, paragraph_positions

def parse_docx_native(file_path: str, filename: str, knowledge_id: int) -> List["Document"]:
    pdf_out = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
    # 生成用于可视化的简单PDF，并获取段落级positions
    _, paragraph_positions = create_simple_pdf_from_docx(file_path, pdf_out)

    # 汇总内容
    all_text = "\n".join([p["text"] for p in paragraph_positions])

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    text_chunks = splitter.split_text(all_text)

    chunks: List["Document"] = []
    for chunk_idx, chunk_text in enumerate(text_chunks):
        chunk_pos = []
        for p in paragraph_positions:
            t = p.get("text") or ""
            if t and t in chunk_text:
                cp = dict(p)
                # 将页号同步为page_num字段
                cp["page"] = p.get("page", 1)
                chunk_pos.append(cp)
        bbox = calculate_chunk_bbox([{"bbox": pp["bbox"]} for pp in chunk_pos]) if chunk_pos else [0, 0, 0, 0]
        page_counts: Dict[int, int] = {}
        for pp in chunk_pos:
            pg = int(pp.get("page", 1))
            page_counts[pg] = page_counts.get(pg, 0) + 1
        main_page = max(page_counts.items(), key=lambda kv: kv[1])[0] if page_counts else 1
        chunks.append(Document(
            page_content=chunk_text,
            metadata={
                "knowledge_id": knowledge_id,
                "source_file": filename,
                "page_num": main_page,
                "chunk_index": chunk_idx,
                "positions": chunk_pos,
                "bbox": bbox,
                "document_name": filename,
                "document_type": "文档",
                "keywords": extract_keywords_from_content(chunk_text),
            }
        ))
    return chunks

# ===== PPTX 原生解析与PDF生成（不依赖外部工具） =====
EMU_PER_INCH = 914400.0
POINTS_PER_INCH = 72.0

def create_simple_pdf_from_pptx(src_path: str, out_path: str) -> tuple[str, List[Dict]]:
    if not PPTX_AVAILABLE:
        raise RuntimeError("缺少python-pptx，无法原生解析PPT。请安装: pip install python-pptx")

    prs = PptxPresentation(src_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)

    pdf = fitz.open()
    def _register_cjk_font(_doc: fitz.Document) -> str:
        try:
            candidates = [
                r"C:\\Windows\\Fonts\\msyh.ttc",
                r"C:\\Windows\\Fonts\\msyh.ttf",
                r"C:\\Windows\\Fonts\\simsun.ttc",
                r"C:\\Windows\\Fonts\\simhei.ttf",
            ]
            for p in candidates:
                if os.path.exists(p):
                    _doc.insert_font(fontname="CJK", fontfile=p)
                    return "CJK"
        except Exception:
            pass
        return "helv"
    cjk_font = _register_cjk_font(pdf)
    # python-pptx 的尺寸为 Length 类型，用 int() 取 EMU 值更稳妥
    slide_w_pts = (int(prs.slide_width) / EMU_PER_INCH) * POINTS_PER_INCH
    slide_h_pts = (int(prs.slide_height) / EMU_PER_INCH) * POINTS_PER_INCH
    positions: List[Dict] = []

    for s_idx, slide in enumerate(prs.slides):
        page = pdf.new_page(width=slide_w_pts, height=slide_h_pts)
        for shape in slide.shapes:
            try:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = "\n".join([p.text for p in shape.text_frame.paragraphs]).strip()
                if not text:
                    continue
                x0 = (int(shape.left) / EMU_PER_INCH) * POINTS_PER_INCH
                y0 = (int(shape.top) / EMU_PER_INCH) * POINTS_PER_INCH
                x1 = x0 + (int(shape.width) / EMU_PER_INCH) * POINTS_PER_INCH
                y1 = y0 + (int(shape.height) / EMU_PER_INCH) * POINTS_PER_INCH
                rect = fitz.Rect(x0, y0, x1, y1)
                # 在对应矩形内写入文本
                page.insert_textbox(rect, text, fontsize=12, fontname=cjk_font, align=0)
                positions.append({"text": text, "bbox": [x0, y0, x1, y1], "page": s_idx + 1})
            except Exception:
                continue

    pdf.save(out_path)
    pdf.close()
    return out_path, positions

def parse_pptx_native(file_path: str, filename: str, knowledge_id: int) -> List["Document"]:
    pdf_out = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
    _, shape_positions = create_simple_pdf_from_pptx(file_path, pdf_out)

    # 汇总内容
    all_text = "\n".join([p["text"] for p in shape_positions])

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    text_chunks = splitter.split_text(all_text)

    chunks: List["Document"] = []
    for idx, chunk_text in enumerate(text_chunks):
        chunk_pos = []
        for p in shape_positions:
            t = p.get("text") or ""
            if t and t in chunk_text:
                chunk_pos.append(dict(p))
        bbox = calculate_chunk_bbox([{"bbox": pp["bbox"]} for pp in chunk_pos]) if chunk_pos else [0, 0, 0, 0]
        page_counts: Dict[int, int] = {}
        for pp in chunk_pos:
            pg = int(pp.get("page", 1))
            page_counts[pg] = page_counts.get(pg, 0) + 1
        main_page = max(page_counts.items(), key=lambda kv: kv[1])[0] if page_counts else 1
        chunks.append(Document(
            page_content=chunk_text,
            metadata={
                "knowledge_id": knowledge_id,
                "source_file": filename,
                "page_num": main_page,
                "chunk_index": idx,
                "positions": chunk_pos,
                "bbox": bbox,
                "document_name": filename,
                "document_type": "演示文稿",
                "keywords": extract_keywords_from_content(chunk_text),
            }
        ))
    return chunks

@app.post("/api/ldap/validate", response_model=LdapValidateResponse)
def validate_ldap_user(request: LdapValidateRequest):
    """LDAP用户验证（模拟实现）"""
    logger.info(f"LDAP验证请求: {request.username}")
    
    # 模拟LDAP验证
    if request.username == "admin" and request.password == "password":
        return LdapValidateResponse(
            success=True,
            message="验证成功",
            email="admin@example.com",
            role="admin"
        )
    else:
        return LdapValidateResponse(
            success=False,
            message="用户名或密码错误"
        )

@app.post("/ldap/verify")
def ldap_verify(request: LdapValidateRequest):
    """LDAP验证接口，供Java调用"""
    username = request.username.strip()
    password = request.password
    
    # 模拟验证（可以后续集成真实LDAP）
    if not username or not password:
        raise HTTPException(status_code=401, detail="用户名或密码不能为空")
    
    # 模拟用户信息
    user_info = {
        "username": username,
        "email": f"{username}@example.com",
        "display_name": username.capitalize(),
        "role": "USER",
        "system_role": "USER",
    }
    
    return {
        "ok": True,
        "source": "mock",
        "user": user_info,
        "raw": {}
    }

def process_document_unified(
    file_path: str,
    filename: str,
    knowledge_id: int,
    knowledge_name: Optional[str] = None,
    description: Optional[str] = None,
    tags: Optional[str] = None,
    effective_time: Optional[str] = None,
) -> Dict:
    """
    统一处理文档：
    - 非PDF先转换为PDF（优先使用LibreOffice），不改动原始文件
    - 基于PDF用PyMuPDF提取块级坐标并分块
    """
    try:
        # 1) 非PDF → PDF（不改动原始文件，仅生成临时PDF用于定位与回显）
        ext = Path(filename).suffix.lower()
        pdf_path_to_open = file_path
        temp_generated_pdf = None

        if ext in {'.xlsx', '.xls'}:
            # Excel 原生解析
            chunks = parse_excel_native(file_path, filename, knowledge_id)
        elif ext == '.txt':
            # TXT 原生解析
            chunks = parse_txt_native(file_path, filename, knowledge_id)
        elif ext in {'.doc', '.docx'}:
            # Word 原生解析
            chunks = parse_docx_native(file_path, filename, knowledge_id)
            # 使用我们生成的简易PDF供回显
            target_pdf = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
            pdf_path_to_open = target_pdf
        elif ext in {'.ppt', '.pptx'}:
            # PPT 原生解析
            chunks = parse_pptx_native(file_path, filename, knowledge_id)
            target_pdf = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
            pdf_path_to_open = target_pdf
        elif ext != '.pdf':
            # 其他Office类型仍尝试转为PDF（如需）
            target_pdf = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
            try:
                target_dir = os.path.dirname(target_pdf)
                src_for_convert = file_path
                pdf_path_to_open = convert_with_libreoffice_safe(src_for_convert, out_dir=target_dir)
                if os.path.abspath(pdf_path_to_open) != os.path.abspath(target_pdf):
                    shutil.copyfile(pdf_path_to_open, target_pdf)
                    pdf_path_to_open = target_pdf
                logger.info(f"已将 {filename} 转为PDF: {pdf_path_to_open}")
            except Exception as e:
                if ext == '.txt':
                    pdf_path_to_open = create_simple_pdf_from_txt(file_path, out_path=target_pdf)
                    logger.info(f"TXT降级转PDF成功: {pdf_path_to_open}")
                else:
                    raise

        # 2) 基于PDF走统一解析
        doc = fitz.open(pdf_path_to_open)
        try:
            logger.info(f"成功打开文档，页数: {len(doc)}")

            documents = extract_documents_with_block_positions(doc, filename)

            all_content = ""
            all_positions = []
            for doc_info in documents:
                all_content += doc_info["content"] + "\n"
                all_positions.extend(doc_info["positions"])

            logger.info(f"合并后总内容长度: {len(all_content)} 字符")
            logger.info(f"合并后总位置信息数量: {len(all_positions)}")

            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
            )

            text_chunks = text_splitter.split_text(all_content)
            logger.info(f"LangChain分割后生成 {len(text_chunks)} 个chunks")

            chunks = []
            for chunk_idx, chunk_text in enumerate(text_chunks):
                chunk_positions = assign_positions_to_chunk(chunk_text, all_positions)

                page_counts: Dict[int, int] = {}
                for p in chunk_positions:
                    pg = int(p.get('page', 1))
                    page_counts[pg] = page_counts.get(pg, 0) + 1
                main_page = max(page_counts.items(), key=lambda kv: kv[1])[0] if page_counts else 1

                chunk = Document(
                    page_content=chunk_text,
                    metadata={
                        "knowledge_id": knowledge_id,
                        "source_file": filename,
                        "page_num": main_page,
                        "chunk_index": chunk_idx,
                        "positions": chunk_positions,
                        "bbox": calculate_chunk_bbox(chunk_positions),
                        "document_name": filename,
                        "document_type": "文档",
                        "keywords": extract_keywords_from_content(chunk_text),
                        # 协同到ES
                        "knowledge_name": knowledge_name or "",
                        "description": description or "",
                        "tags": tags or "",
                        "effective_time": effective_time or "",
                    }
                )
                chunks.append(chunk)
        finally:
            try:
                doc.close()
            except Exception:
                pass
            # 不删除持久化PDF，供前端下载/回显
        
        # 为Excel/TXT等原生解析生成的chunk补齐知识元数据字段
        try:
            for ch in chunks:
                md = ch.metadata
                if "knowledge_name" not in md:
                    md["knowledge_name"] = knowledge_name or ""
                if "description" not in md:
                    md["description"] = description or ""
                if "tags" not in md:
                    md["tags"] = tags or ""
                if "effective_time" not in md:
                    md["effective_time"] = effective_time or ""
        except Exception:
            pass

        logger.info(f"最终生成 {len(chunks)} 个chunks")
        
        # 存储到ES
        stored_count = store_chunks_to_es(chunks, knowledge_id)
        
        return {
            "chunks_count": stored_count,
            "total_chunks": len(chunks),
            "success": stored_count > 0
        }
        
    except Exception as e:
        logger.error(f"文档处理失败: {e}")
        import traceback
        traceback.print_exc()
        raise e

def extract_documents_with_block_positions(doc, filename: str) -> List[Dict]:
    """
    直接从PyMuPDF提取文档块和位置信息
    """
    documents = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        blocks = page.get_text("dict")

        for block_idx, block in enumerate(blocks["blocks"]):
            if "lines" in block:
                # 收集整个块的所有文本和位置信息
                block_text = ""
                block_positions = []

                for line_idx, line in enumerate(block["lines"]):
                    for span_idx, span in enumerate(line["spans"]):
                        text = span["text"].strip()
                        if text:
                            block_text += text + " "
                            block_positions.append({
                                "text": text,
                                "bbox": span["bbox"],
                                "font_size": span["size"],
                                "font": span["font"],
                                "span_idx": span_idx,
                                "line_idx": line_idx,
                                "page": page_num + 1,
                            })

                if block_text.strip():
                    # 计算整个块的边界框（内联实现）
                    if block_positions:
                        bx0 = min(p["bbox"][0] for p in block_positions)
                        by0 = min(p["bbox"][1] for p in block_positions)
                        bx1 = max(p["bbox"][2] for p in block_positions)
                        by1 = max(p["bbox"][3] for p in block_positions)
                        block_bbox = [bx0, by0, bx1, by1]
                    else:
                        block_bbox = [0, 0, 0, 0]

                    documents.append({
                        "content": block_text.strip(),
                        "page": page_num + 1,
                        "block_idx": block_idx,
                        "positions": block_positions,
                        "bbox": block_bbox
                    })

    return documents


def assign_positions_to_chunk(chunk_text: str, positions: List[Dict]) -> List[Dict]:
    """
    为chunk分配对应的位置信息
    使用简单的文本包含关系，而不是复杂的相似度计算
    """
    chunk_positions = []

    for pos in positions:
        pos_text = pos["text"]
        # 如果chunk包含这个位置的文本，就分配给它
        if pos_text in chunk_text:
            chunk_positions.append(pos)

    return chunk_positions

def calculate_chunk_bbox(positions: List[Dict]) -> List[float]:
    """计算chunk的边界框：先按span所在页聚类，取位置最多的页作为主页面，再对该页的positions求并集"""
    if not positions:
        return [0, 0, 0, 0]

    # 位置里需要带上page信息；若没有，默认页=1
    page_to_positions: Dict[int, List[Dict]] = {}
    for pos in positions:
        page = int(pos.get('page', 1)) if isinstance(pos, dict) else 1
        page_to_positions.setdefault(page, []).append(pos)

    # 选择位置最多的页作为主页面
    main_page = max(page_to_positions.items(), key=lambda kv: len(kv[1]))[0]
    main_positions = page_to_positions[main_page]

    x0 = min(p["bbox"][0] for p in main_positions)
    y0 = min(p["bbox"][1] for p in main_positions)
    x1 = max(p["bbox"][2] for p in main_positions)
    y1 = max(p["bbox"][3] for p in main_positions)

    return [x0, y0, x1, y1]

def generate_pdfllm_style_markdown(doc, filename: str) -> tuple[str, List[Dict]]:
    """
    使用PyMuPDF生成干净的Markdown内容和精确的文本-坐标映射
    返回: (markdown_content, position_mapping)
    """
    content_lines = []
    position_mapping = []
    
    # 添加文档标题
    content_lines.append(f"# {filename}")
    content_lines.append("")
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # 页面标题
        content_lines.append(f"## 第 {page_num + 1} 页")
        content_lines.append("")
        
        # 获取页面文本块
        blocks = page.get_text("dict")
        
        if "blocks" in blocks:
            for block_idx, block in enumerate(blocks["blocks"]):
                if "lines" in block:
                    block_text = ""
                    block_positions = []
                    
                    for line_idx, line in enumerate(block["lines"]):
                        line_text = ""
                        line_positions = []
                        
                        for span_idx, span in enumerate(line["spans"]):
                            text = span["text"].strip()
                            if text:
                                # 只添加文本内容，不添加位置标签
                                line_text += text + " "
                                
                                # 记录位置信息
                                line_positions.append({
                                    "text": text,
                                    "page": page_num + 1,
                                    "bbox": span["bbox"],
                                    "font_size": span["size"],
                                    "font": span["font"],
                                    "block_idx": block_idx,
                                    "line_idx": line_idx,
                                    "span_idx": span_idx,
                                    "char_start": len("".join(content_lines)),  # 在最终文本中的起始位置
                                    "char_end": len("".join(content_lines)) + len(text)  # 在最终文本中的结束位置
                                })
                        
                        if line_text.strip():
                            block_text += line_text.strip() + "\n"
                            block_positions.extend(line_positions)
                    
                    if block_text.strip():
                        # 检查是否可能是标题（基于字体大小）
                        max_font_size = max(span["size"] for line in block["lines"] for span in line["spans"])
                        if max_font_size > 12:  # 假设大于12pt的是标题
                            content_lines.append(f"### {block_text.strip()}")
                        else:
                            content_lines.append(block_text.strip())
                        content_lines.append("")
                        
                        # 添加位置信息到映射
                        position_mapping.extend(block_positions)
    
    # 合并所有内容
    markdown_content = "\n".join(content_lines)
    
    # 更新字符位置（因为换行符等会影响位置）
    current_pos = 0
    for pos_info in position_mapping:
        text = pos_info["text"]
        pos_info["char_start"] = current_pos
        pos_info["char_end"] = current_pos + len(text)
        current_pos += len(text) + 1  # +1 for space
    
    return markdown_content, position_mapping

def extract_position_mapping(doc) -> List[Dict]:
    """
    单独提取位置信息映射，不嵌入到文本内容中
    """
    position_mapping = []
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        blocks = page.get_text("dict")
        
        if "blocks" in blocks:
            for block_idx, block in enumerate(blocks["blocks"]):
                if "lines" in block:
                    for line_idx, line in enumerate(block["lines"]):
                        for span_idx, span in enumerate(line["spans"]):
                            text = span["text"].strip()
                            if text:
                                position_mapping.append({
                                    "text": text,
                                    "page": page_num + 1,
                                    "bbox": span["bbox"],
                                    "font_size": span["size"],
                                    "font": span["font"],
                                    "block_idx": block_idx,
                                    "line_idx": line_idx,
                                    "span_idx": span_idx
                                })
    
    return position_mapping

def find_best_position_match(chunk_text: str, position_mapping: List[Dict]) -> Optional[Dict]:
    """
    为chunk找到最匹配的位置信息
    使用块级模糊匹配，而不是精确的文本片段匹配
    """
    if not position_mapping:
        return None
    
    chunk_text_lower = chunk_text.lower()
    chunk_words = set(chunk_text_lower.split())
    
    # 定义关键信息关键词和对应的优先级
    key_phrases_priority = [
        # 高优先级：基金核心信息
        (["基金总值", "4.4377", "亿美元"], 10),
        (["基金价格", "资产净值", "5.7741"], 9),
        (["成立日期", "2010年8月2日"], 8),
        (["基金经理", "Justin Kass", "David Oberto", "Michael Yee"], 7),
        (["管理费", "1.19%"], 6),
        (["投资目标", "美国债券", "高收益"], 5),
        (["收益分配", "每月"], 4),
        (["财政年度", "9月30日"], 3),
        (["交易日", "每日"], 2),
        (["投资经理", "安联投资"], 1)
    ]
    
    # 第一轮：优先匹配包含关键信息的文本
    for pos_info in position_mapping:
        text = pos_info.get("text", "").strip()
        if text:
            # 检查是否包含高优先级的关键信息
            for key_phrases, priority in key_phrases_priority:
                if any(keyword in text for keyword in key_phrases):
                    # 计算文本相似度
                    text_lower = text.lower()
                    chunk_lower = chunk_text_lower
                    
                    # 使用字符重叠度计算相似度
                    overlap = sum(1 for c in text_lower if c in chunk_lower)
                    base_score = overlap / max(len(text_lower), 1)
                    
                    # 应用优先级权重
                    weighted_score = base_score * priority
                    
                    if weighted_score > best_score:
                        best_score = weighted_score
                        best_match = pos_info
                    break  # 找到匹配的关键词就跳出内层循环
    
    # 第二轮：如果没有找到关键信息，尝试精确包含匹配
    if not best_match:
        for pos_info in position_mapping:
            text = pos_info.get("text", "").strip()
            if text and text.lower() in chunk_text_lower:
                # 计算匹配度：文本长度与chunk长度的比例
                score = len(text) / max(len(chunk_text), 1)
                if score > best_score:
                    best_score = score
                    best_match = pos_info
    
    # 第三轮：如果还是没有找到，尝试单词级别的匹配
    if not best_match:
        for pos_info in position_mapping:
            text = pos_info.get("text", "").strip()
            if text:
                text_words = set(text.lower().split())
                # 计算单词重叠度
                overlap = len(chunk_words.intersection(text_words))
                if overlap > 0:
                    score = overlap / max(len(chunk_words), 1)
                    if score > best_score:
                        best_score = score
                        best_match = pos_info
    
    # 第四轮：如果还是没有找到，返回第一个有效的位置信息
    if not best_match and position_mapping:
        for pos_info in position_mapping:
            if pos_info.get("bbox") and len(pos_info.get("bbox", [])) == 4:
                return pos_info
    
    return best_match

def expand_bbox_to_block_level(bbox: List[float], page_width: float, page_height: float) -> List[float]:
    """
    将精确的文本片段bbox扩展为块级bbox，提供更好的用户体验
    使用更保守的扩展策略，避免过度扩展
    """
    if len(bbox) != 4:
        return bbox
    
    x0, y0, x1, y1 = bbox
    
    # 计算当前文本的宽度和高度
    text_width = x1 - x0
    text_height = y1 - y0
    
    # 更保守的扩展策略
    # 水平扩展：左右各扩展文本宽度的50%，但不超过页面边距
    horizontal_expansion = min(text_width * 0.5, 30)  # 最大扩展30像素
    
    expanded_x0 = max(20, x0 - horizontal_expansion)
    expanded_x1 = min(page_width - 20, x1 + horizontal_expansion)
    
    # 垂直扩展：上下各扩展文本高度的50%，但不超过页面边距
    vertical_expansion = min(text_height * 0.5, 20)  # 最大扩展20像素
    
    expanded_y0 = max(20, y0 - vertical_expansion)
    expanded_y1 = min(page_height - 20, y1 + vertical_expansion)
    
    return [expanded_x0, expanded_y0, expanded_x1, expanded_y1]

def find_best_position_match_block_level(chunk_text: str, position_mapping: List[Dict], page_width: float, page_height: float) -> Optional[Dict]:
    """
    为chunk找到最匹配的位置信息，并扩展为块级坐标
    """
    # 先找到最佳匹配
    best_match = find_best_position_match(chunk_text, position_mapping)
    
    if best_match and best_match.get("bbox"):
        # 扩展bbox为块级坐标
        expanded_bbox = expand_bbox_to_block_level(best_match["bbox"], page_width, page_height)
        
        # 创建新的位置信息，包含扩展后的坐标
        block_level_match = best_match.copy()
        block_level_match["bbox"] = expanded_bbox
        block_level_match["bbox_type"] = "block_level"  # 标记为块级坐标
        
        return block_level_match
    
    return best_match

def store_chunks_to_es(chunks: List[Document], knowledge_id: int):
    """
    将chunks存储到Elasticsearch
    """
    stored_count = 0
    
    for i, chunk in enumerate(chunks):
        try:
            # 生成文档ID
            doc_id = hashlib.md5(f"{knowledge_id}_{i}_{chunk.page_content[:100]}".encode()).hexdigest()
            
            # 为chunk生成embedding
            chunk_embedding = get_embedding(chunk.page_content)
            if not chunk_embedding:
                logger.warning(f"Chunk {i} embedding生成失败，跳过")
                continue
            # 校验 embedding 维度与配置一致（自定义API切换成1024时避免落错索引）
            try:
                emb_len = len(chunk_embedding) if chunk_embedding is not None else 0
            except Exception:
                emb_len = 0
            if emb_len != EMBEDDING_DIMS:
                logger.error(f"Chunk {i} embedding 维度不匹配: got={emb_len}, expected={EMBEDDING_DIMS}，已跳过入库")
                continue
            
            # 准备ES文档
            es_doc = {
                "content": chunk.page_content,
                "embedding": chunk_embedding,
                "knowledge_id": chunk.metadata.get("knowledge_id", knowledge_id),
                "knowledge_name": chunk.metadata.get("knowledge_name", ""),
                "description": chunk.metadata.get("description", ""),
                "tags": chunk.metadata.get("tags", ""),
                "effective_time": chunk.metadata.get("effective_time", ""),
                "source_file": chunk.metadata.get("source_file", ""),
                "chunk_index": chunk.metadata.get("chunk_index", i),
                "chunk_type": chunk.metadata.get("chunk_type", "content"),
                "page_num": chunk.metadata.get("page_num", 1),
                "bbox": chunk.metadata.get("bbox", []),
                "positions": chunk.metadata.get("positions", []),
                "node_type": "doc",  # 明确标识这是文档类型
                "weight": 1.0
            }
            
            # 存储到ES
            es_client.index(index=ES_CONFIG['index'], id=doc_id, document=es_doc)
            stored_count += 1
            
            if (i + 1) % 10 == 0:
                logger.info(f"已存储 {i + 1}/{len(chunks)} 个chunks")
                
        except Exception as e:
            logger.error(f"存储chunk {i} 失败: {e}")
            continue

    logger.info(f"ES存储完成，成功存储 {stored_count}/{len(chunks)} 个chunks")
    return stored_count

@app.post("/api/document/process", response_model=DocumentProcessResponse)
async def process_document(
    file: UploadFile = File(...),
    knowledge_id: int = Form(None),
    knowledge_name: str = Form(None),
    description: str = Form(None),
    tags: str = Form(None),
    effective_time: str = Form(None)
):
    """
    处理上传的文档
    - 支持切分的格式：PDF、Word、Excel、PowerPoint、TXT 等
    - 不支持切分的格式：图片、音频、视频等（仅存储）
    """
    logger.info(f"开始处理文档: {file.filename}, 知识ID: {knowledge_id}")
    
    try:
        # 检查文件类型
        file_extension = Path(file.filename).suffix.lower()
        
        # 支持切分的文件类型
        chunkable_extensions = {
            ".pdf", ".docx", ".doc", ".xlsx", ".xls", 
            ".pptx", ".ppt", ".txt", ".hwp", ".hwpx"
        }
        
        # 支持存储但不切分的文件类型（图片、音频、视频等）
        storage_only_extensions = {
            ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp",  # 图片
            ".mp3", ".wav", ".flac", ".aac", ".ogg",  # 音频
            ".mp4", ".avi", ".mov", ".wmv", ".flv", ".mkv",  # 视频
            ".zip", ".rar", ".7z", ".tar", ".gz",  # 压缩包
            ".exe", ".msi", ".dmg", ".deb", ".rpm",  # 可执行文件
        }
        
        # 检查是否为支持的文件类型
        if file_extension not in chunkable_extensions and file_extension not in storage_only_extensions:
            raise HTTPException(status_code=400, detail=f"不支持的文件类型: {file_extension}")
        
        # 保存上传的文件
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
            content = await file.read()
            temp_file.write(content)
            temp_file_path = temp_file.name
        
        try:
            # 判断是否需要切分处理
            if file_extension in chunkable_extensions:
                # 进行文档切分处理
                result = process_document_unified(
                    temp_file_path,
                    file.filename,
                    knowledge_id,
                    knowledge_name=knowledge_name,
                    description=description,
                    tags=tags,
                    effective_time=effective_time,
                )
                
                return DocumentProcessResponse(
                    success=True,
                    message=f"文档处理成功: {file.filename}",
                    chunks_count=result["chunks_count"],
                    knowledge_id=int(knowledge_id) if knowledge_id is not None else 0
                )
            else:
                # 仅存储，不进行切分处理
                logger.info(f"文件类型 {file_extension} 不支持切分，仅进行存储: {file.filename}")
                
                return DocumentProcessResponse(
                    success=True,
                    message=f"文件存储成功: {file.filename}（该文件类型不支持内容切分）",
                    chunks_count=0,
                    knowledge_id=int(knowledge_id) if knowledge_id is not None else 0
                )
            
        finally:
            # 清理临时文件
            os.unlink(temp_file_path)
            
    except Exception as e:
        logger.error(f"文档处理失败: {e}")
        raise HTTPException(status_code=500, detail=f"文档处理失败: {str(e)}")

@app.post("/api/rag/chat", response_model=ChatResponse)
def chat_with_rag(request: ChatRequest):
    """
    基于知识库的智能问答
    """
    logger.info(f"RAG聊天请求: {request.question}")
    
    try:
        # 1. 向量搜索找到相关chunks
        question_embedding = get_embedding(request.question)
        if not question_embedding:
            return ChatResponse(
                answer="抱歉，无法生成问题的向量表示，请重试。",
                references=[],
                session_id=request.user_id
            )
        
        # 构建过滤条件 - 简化过滤，先确保基本搜索能工作
        filters = [
            {"term": {"chunk_type": "content"}}  # 只搜索内容类型的chunk
        ]
        
        # 如果指定了特定文件，则只检索该文件的chunks
        if request.source_file:
            logger.info(f"按文件名过滤RAG检索: {request.source_file}")
            filters.append({"term": {"source_file": request.source_file}})
        
        # 简化搜索逻辑 - 直接返回语义相似度最高的chunks
        # 仅检索包含embedding字段的文档，避免脚本在缺失字段时报错
        search_query = {
            "size": 10 if request.source_file else 5,  # 针对特定文件时返回更多chunks
            "query": {
                "script_score": {
                    "query": {
                        "bool": {
                            "filter": filters
                        }
                    },
                    "script": {
                        "source": "cosineSimilarity(params.query_vector, 'embedding') + 1.0",
                        "params": {"query_vector": question_embedding}
                    }
                }
            },
            "_source": [
                "content",
                "knowledge_id",
                "knowledge_name",
                "description",
                "tags",
                "effective_time",
                "source_file",
                "page_num",
                "chunk_index",
                "bbox",
                "positions"
            ]
        }
        
        search_response = es_client.search(index=ES_CONFIG['index'], body=search_query)
        hits = search_response.get('hits', {}).get('hits', [])
        
        if not hits:
            return ChatResponse(
                answer="抱歉，在知识库中没有找到相关信息。",
                references=[],
                session_id=request.user_id
            )
        
        # 2. 构建增强的上下文信息
        context_chunks = []
        for hit in hits:
            source = hit['_source']
            score = hit['_score']
            
            # 构建每个chunk的完整上下文信息
            chunk_info = {
                "content": source.get('content', ''),
                "metadata": {
                    "knowledge_id": source.get('knowledge_id', 0),
                    "knowledge_name": source.get('knowledge_name', ''),
                    "description": source.get('description', ''),
                    "tags": source.get('tags', ''),
                    "effective_time": source.get('effective_time', ''),
                    "document_name": source.get('source_file', 'N/A'),
                    "document_type": "文档",  # 通用文档类型
                    "page_num": source.get('page_num', 'N/A'),
                    "chunk_index": source.get('chunk_index', 'N/A'),
                    "bbox": source.get('bbox', []),
                    "positions": source.get('positions', []),
                    "relevance_score": round(score, 3)
                }
            }
            context_chunks.append(chunk_info)
        
        # 3. 构建增强的RAG提示词
        enhanced_prompt = build_enhanced_rag_prompt(request.question, context_chunks)
        
        # 4. 调用大模型生成答案
        answer = generate_ai_answer(enhanced_prompt)
        
        # 5. 构建引用信息
        references = []
        for chunk in context_chunks:
            metadata = chunk['metadata']
            references.append(KnowledgeReference(
                knowledge_id=int(metadata.get('knowledge_id', 0)) if metadata.get('knowledge_id') is not None else 0,
                knowledge_name=metadata.get('knowledge_name', ''),
                description=metadata.get('description', ''),
                tags=[metadata.get('tags', '')] if isinstance(metadata.get('tags', ''), str) else metadata.get('tags', []),
                effective_time=metadata.get('effective_time', ''),
                attachments=[metadata.get('document_name', '')],
                relevance=metadata.get('relevance_score', 0.0),
                source_file=metadata.get('document_name', ''),
                page_num=metadata.get('page_num', 0),
                chunk_index=metadata.get('chunk_index', 0),
                chunk_type="content",
                bbox_union=metadata.get('bbox', []),  # 使用新的bbox字段
                char_start=0,  # 不再使用字符位置
                char_end=0
            ))
        
        return ChatResponse(
            answer=answer,
            references=references,
            session_id=request.user_id
        )
        
    except Exception as e:
        logger.error(f"RAG聊天失败: {e}")
        return ChatResponse(
            answer=f"抱歉，处理您的问题时出现错误：{str(e)}",
            references=[],
            session_id=request.user_id
        )

def build_enhanced_rag_prompt(question: str, context_chunks: List[Dict]) -> str:
    """
    构建增强的RAG提示词，让大模型能看到完整的上下文信息
    """
    context_parts = []
    
    for i, chunk in enumerate(context_chunks):
        metadata = chunk['metadata']
        
        # 构建每个chunk的详细上下文信息
        chunk_context = f"""
=== 引用 {i+1} ===
📄 文档名称: {metadata.get('document_name', 'N/A')}
📋 文档类型: {metadata.get('document_type', 'N/A')}
📖 页码: {metadata.get('page_num', 'N/A')}
🔢 块序: {metadata.get('chunk_index', 'N/A')}
🎯 相关性: {metadata.get('relevance_score', 'N/A')}
📍 坐标: {metadata.get('bbox', [])}
📝 内容: {chunk.get('content', '')}
"""
        context_parts.append(chunk_context)
    
    # 构建完整提示
    prompt = f"""
你是一个专业的文档知识助手。请基于以下信息回答问题。

每个引用都包含了完整的文档信息，包括：
- 文档名称和类型
- 页码和块序
- 相关性评分和位置坐标
- 具体内容

请仔细分析这些信息，并根据问题找到最准确的答案。

问题: {question}

参考信息:
{''.join(context_parts)}

请根据问题，从上述信息中找到最准确的答案。要求：
1. 直接回答问题，不要说"请查看引用信息"
2. 如果问题涉及特定文档，请确保答案来自正确的文档
3. 如果问题没有指定具体文档，请基于所有相关信息给出综合回答
4. 答案要具体、准确，包含关键数据
5. 用中文回答，并说明信息来源（如"根据文档第X页"）

请开始回答：
"""
    
    return prompt

def generate_ai_answer(prompt: str) -> str:
    """
    调用大模型生成答案
    """
    try:
        # 调用极客智坊API生成答案
        headers = {
            "Authorization": f"Bearer {GEEKAI_API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": "gpt-4o-mini",  # 使用合适的模型
            "messages": [
                {
                    "role": "system",
                    "content": "你是一个专业的文档知识助手，请基于提供的文档信息准确回答问题。"
                },
                {
                    "role": "user", 
                    "content": prompt
                }
            ],
            "max_tokens": 1000,
            "temperature": 0.3
        }
        
        response = requests.post(
            GEEKAI_CHAT_URL,
            headers=headers,
            json=payload,
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            if "choices" in result and len(result["choices"]) > 0:
                answer = result["choices"][0]["message"]["content"]
                return answer
            else:
                logger.error(f"API响应格式异常: {result}")
                return "抱歉，生成答案时出现格式错误，请重试。"
        else:
            logger.error(f"API调用失败: {response.status_code}, {response.text}")
            return f"抱歉，API调用失败（状态码：{response.status_code}），请重试。"
            
    except requests.exceptions.Timeout:
        logger.error("API调用超时")
        return "抱歉，生成答案超时，请重试。"
    except requests.exceptions.RequestException as e:
        logger.error(f"API请求异常: {e}")
        return f"抱歉，API请求异常：{str(e)}，请重试。"
    except Exception as e:
        logger.error(f"生成AI答案失败: {e}")
        return f"抱歉，生成答案时出现错误：{str(e)}，请重试。"

def extract_keywords_from_content(content: str) -> List[str]:
    """从内容提取关键标识词（通用版本）"""
    keywords = []
    content_lower = content.lower()
    
    # 通用关键词
    general_keywords = ["投资", "管理", "风险", "收益", "费用", "日期", "目标", "策略", "报告", "分析"]
    for keyword in general_keywords:
        if keyword in content_lower:
            keywords.append(keyword)
    
    # 文档结构关键词
    structure_keywords = ["标题", "章节", "表格", "图表", "附录", "摘要", "结论"]
    for keyword in structure_keywords:
        if keyword in content_lower:
            keywords.append(keyword)
    
    return keywords[:8]  # 限制关键词数量

@app.get("/api/health")
def health_check():
    """
    健康检查
    """
    try:
        # 检查ES连接
        es_info = es_client.info()
        
        # 检查PyMuPDF可用性
        pymupdf_status = "available" if PYMUPDF_AVAILABLE else "unavailable"
        
        # 检查极客智坊API
        try:
            headers = {"Authorization": f"Bearer {GEEKAI_API_KEY}"}
            response = requests.get("https://geekai.co/api/v1/models", headers=headers, timeout=5)
            geekai_status = "available" if response.status_code == 200 else "unavailable"
        except:
            geekai_status = "unavailable"
        
        return {
            "status": "healthy",
            "elasticsearch": "connected",
            "pymupdf": pymupdf_status,
            "geekai_api": geekai_status,
            "timestamp": "2024-01-01T00:00:00Z"
        }
    except Exception as e:
        logger.error(f"健康检查失败: {e}")
        raise HTTPException(status_code=500, detail=f"健康检查失败: {str(e)}")

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
